#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
UltimateReader.py
# 专业级通用文档阅读器（PyQt5重构版）
作者：LYP
GitHub：https://github.com/lyp0746
邮箱：1610369302@qq.com
版本：4.0.0
"""

import sys
import json
import time
import os
from pathlib import Path
from collections import defaultdict
from typing import Optional, List, Dict, Tuple
from datetime import datetime

from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, QToolBar,
    QPushButton, QLabel, QFileDialog, QScrollArea, QSplitter, QListWidget,
    QLineEdit, QComboBox, QStatusBar, QMessageBox, QAction, QTreeWidget,
    QTreeWidgetItem, QTabWidget, QTextEdit, QListWidgetItem, QInputDialog,
    QMenu, QMenuBar, QDialog, QDialogButtonBox, QProgressDialog, QProgressBar,
    QTextBrowser, QGroupBox, QCheckBox, QSpinBox, QDoubleSpinBox, QRadioButton,
    QButtonGroup, QSlider, QTableWidget, QTableWidgetItem, QHeaderView,
    QFrame, QGridLayout, QFormLayout
)
from PyQt5.QtCore import (
    Qt, QThread, pyqtSignal, QTimer, QSize, QRect, QPoint, QPointF,
    QRectF, QSettings, QByteArray, QMimeData, QUrl, QEasingCurve,
    QPropertyAnimation, QParallelAnimationGroup, pyqtSlot
)
from PyQt5.QtGui import (
    QPixmap, QImage, QPalette, QColor, QIcon, QPainter, QKeySequence,
    QFont, QTextCursor, QTextCharFormat, QBrush, QPen, QDragEnterEvent,
    QDropEvent, QWheelEvent, QPaintEvent, QMouseEvent, QCursor, QTextDocument,
    QTextFormat, QFontMetrics
)
from PyQt5.QtPrintSupport import QPrinter, QPrintDialog, QPrintPreviewDialog

import fitz  # PyMuPDF

# ---------- 可选依赖 ----------
try:
    from docx import Document

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from odf import text as odf_text, teletype
    from odf.opendocument import load as odf_load

    ODT_AVAILABLE = True
except ImportError:
    ODT_AVAILABLE = False

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image

    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# ---------- 全局配置 ----------
CONFIG_FILE = Path.cwd() / ".ultimate_reader" / "config.json"
CONFIG_FILE.parent.mkdir(parents=True, exist_ok=True)


# ==================== 配置管理 ====================
class ConfigManager:
    """集中配置管理器"""

    @staticmethod
    def load() -> dict:
        if CONFIG_FILE.exists():
            try:
                with open(CONFIG_FILE, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except Exception as e:
                print(f"加载配置失败: {e}")
        return ConfigManager.default_config()

    @staticmethod
    def save(config: dict):
        try:
            with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"保存配置失败: {e}")

    @staticmethod
    def default_config() -> dict:
        return {
            "window": {"x": 100, "y": 100, "width": 1600, "height": 900},
            "recent_files": [],
            "documents": {},
            "reading_positions": {},
            "theme": "light",
            "auto_save": True,
            "auto_save_interval": 300,  # 秒
            "smooth_scroll": True,
            "preload_pages": 3,
            "ocr_language": "chi_sim+eng",
            "last_opened_dir": str(Path.cwd())
        }

    @staticmethod
    def get_document_state(file_path: str) -> dict:
        config = ConfigManager.load()
        return config.get("documents", {}).get(file_path, {})

    @staticmethod
    def set_document_state(file_path: str, state: dict):
        config = ConfigManager.load()
        docs = config.get("documents", {})
        docs[file_path] = state
        config["documents"] = docs
        ConfigManager.save(config)

    # ==================== 渲染线程 ====================


class DocumentRenderThread(QThread):
    """异步文档渲染线程"""
    page_rendered = pyqtSignal(int, QPixmap)
    render_failed = pyqtSignal(int, str)

    def __init__(self, doc, page_num: int, zoom: float, rotation: int,
                 dark_mode: bool, high_quality: bool = True):
        super().__init__()
        self.doc = doc
        self.page_num = page_num
        self.zoom = zoom
        self.rotation = rotation
        self.dark_mode = dark_mode
        self.high_quality = high_quality
        self._stop = False

    def stop(self):
        self._stop = True

    def run(self):
        if self._stop:
            return

        try:
            page = self.doc[self.page_num]
            mat = fitz.Matrix(self.zoom / 100.0, self.zoom / 100.0).prerotate(self.rotation)

            # 高质量渲染
            pix = page.get_pixmap(
                matrix=mat,
                alpha=False,
                colorspace=fitz.csRGB,
                clip=None
            )

            img = QImage(
                pix.samples, pix.width, pix.height,
                pix.stride, QImage.Format_RGB888
            )

            if self.dark_mode:
                img.invertPixels()

            if not self._stop:
                pixmap = QPixmap.fromImage(img)
                self.page_rendered.emit(self.page_num, pixmap)

        except Exception as e:
            if not self._stop:
                self.render_failed.emit(self.page_num, str(e))

            # ==================== 页面组件 ====================


class PageWidget(QLabel):
    """图像页面显示组件（支持注释绘制）"""

    annotation_created = pyqtSignal(int, QRect, str)  # page_num, rect, type

    def __init__(self, view, page_num: int):
        super().__init__()
        self.view = view
        self.page_num = page_num
        self.setAlignment(Qt.AlignCenter)
        self.setStyleSheet("background-color: #E8E8E8; margin: 8px;")
        self.setScaledContents(False)

        # 选择区域
        self.selecting = False
        self.selection_start: Optional[QPoint] = None
        self.selection_end: Optional[QPoint] = None

        # 启用鼠标追踪
        self.setMouseTracking(True)

    def mousePressEvent(self, event: QMouseEvent):
        if event.button() == Qt.LeftButton and self.view.annotation_mode:
            self.selecting = True
            self.selection_start = event.pos()
            self.selection_end = event.pos()
            self.update()
        else:
            super().mousePressEvent(event)

    def mouseMoveEvent(self, event: QMouseEvent):
        if self.selecting:
            self.selection_end = event.pos()
            self.update()
        else:
            super().mouseMoveEvent(event)

    def mouseReleaseEvent(self, event: QMouseEvent):
        if self.selecting and event.button() == Qt.LeftButton:
            self.selecting = False
            self.selection_end = event.pos()
            rect = self._get_selection_rect()
            self.selection_start = None
            self.selection_end = None
            self.update()

            if rect and rect.width() > 10 and rect.height() > 10:
                self.annotation_created.emit(
                    self.page_num, rect, self.view.annotation_mode
                )
        else:
            super().mouseReleaseEvent(event)

    def _get_selection_rect(self) -> Optional[QRect]:
        if self.selection_start and self.selection_end:
            return QRect(
                min(self.selection_start.x(), self.selection_end.x()),
                min(self.selection_start.y(), self.selection_end.y()),
                abs(self.selection_end.x() - self.selection_start.x()),
                abs(self.selection_end.y() - self.selection_start.y())
            )
        return None

    def paintEvent(self, event: QPaintEvent):
        super().paintEvent(event)

        if not self.pixmap() or self.pixmap().isNull():
            return

        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        painter.setRenderHint(QPainter.SmoothPixmapTransform)

        zoom = self.view.zoom_level / 100.0

        # 绘制搜索高亮
        self._draw_search_highlights(painter, zoom)

        # 绘制注释
        self._draw_annotations(painter, zoom)

        # 绘制当前选择框
        if self.selecting and self.selection_start and self.selection_end:
            rect = self._get_selection_rect()
            if rect:
                pen = QPen(QColor(0, 120, 215), 2, Qt.DashLine)
                painter.setPen(pen)
                painter.setBrush(QColor(0, 120, 215, 30))
                painter.drawRect(rect)

        painter.end()

    def _draw_search_highlights(self, painter: QPainter, zoom: float):
        """绘制搜索结果高亮"""
        search_rects = self.view.search_results_pdf.get(self.page_num, [])
        current_hit = self.view.current_search_hit

        for idx, rect in enumerate(search_rects):
            qrect = QRect(
                int(rect.x0 * zoom), int(rect.y0 * zoom),
                int((rect.x1 - rect.x0) * zoom), int((rect.y1 - rect.y0) * zoom)
            )

            # 当前匹配用不同颜色
            if current_hit and current_hit == (self.page_num, idx):
                color = QColor(255, 140, 0, 180)
            else:
                color = QColor(255, 255, 0, 120)

            painter.fillRect(qrect, color)

    def _draw_annotations(self, painter: QPainter, zoom: float):
        """绘制注释标记"""
        annotations = self.view.annotations_by_page.get(self.page_num, [])

        for ann in annotations:
            rect = ann.get("rect", [0, 0, 0, 0])
            qrect = QRect(
                int(rect[0] * zoom), int(rect[1] * zoom),
                int((rect[2] - rect[0]) * zoom), int((rect[3] - rect[1]) * zoom)
            )

            ann_type = ann.get("type", "highlight")

            if ann_type == "highlight":
                painter.fillRect(qrect, QColor(255, 255, 0, 100))
            elif ann_type == "underline":
                pen = QPen(QColor(0, 100, 255), 2)
                painter.setPen(pen)
                painter.drawLine(qrect.bottomLeft(), qrect.bottomRight())
            elif ann_type == "note":
                # 绘制批注图标
                painter.setBrush(QColor(255, 200, 0, 220))
                painter.setPen(QPen(QColor(200, 150, 0), 1))
                icon_rect = QRect(qrect.left(), qrect.top(), 24, 24)
                painter.drawEllipse(icon_rect)
                painter.setPen(QColor(100, 50, 0))
                font = painter.font()
                font.setPointSize(10)
                font.setBold(True)
                painter.setFont(font)
                painter.drawText(icon_rect, Qt.AlignCenter, "📝")


class TextPageWidget(QTextEdit):
    """文本页面显示组件"""

    def __init__(self, view, page_num: int, editable: bool = False):
        super().__init__()
        self.view = view
        self.page_num = page_num
        self.setReadOnly(not editable)
        self._apply_theme()

        font = QFont("微软雅黑", 11)
        self.setFont(font)

        self.search_selections: List[QTextEdit.ExtraSelection] = []

    def _apply_theme(self):
        theme = self.view.main.current_theme
        if theme == "dark":
            self.setStyleSheet("""  
                background-color: #1E1E1E;   
                color: #D4D4D4;   
                margin: 8px;   
                padding: 20px;  
                border: 1px solid #3C3C3C;  
            """)
        elif theme == "eye_care":
            self.setStyleSheet("""  
                background-color: #C7EDCC;   
                color: #2F4F2F;   
                margin: 8px;   
                padding: 20px;  
                border: 1px solid #8FBC8F;  
            """)
        else:
            self.setStyleSheet("""  
                background-color: white;   
                color: black;   
                margin: 8px;   
                padding: 20px;  
                border: 1px solid #CCCCCC;  
            """)

    def clear_search_highlight(self):
        self.search_selections = []
        self.setExtraSelections([])

    def highlight_all(self, term: str, color: QColor) -> int:
        """高亮所有匹配项"""
        self.clear_search_highlight()
        if not term:
            return 0

        doc = self.document()
        cursor = QTextCursor(doc)
        fmt = QTextCharFormat()
        fmt.setBackground(color)

        count = 0
        selections = []

        while True:
            cursor = doc.find(term, cursor)
            if cursor.isNull():
                break

            sel = QTextEdit.ExtraSelection()
            sel.cursor = cursor
            sel.format = fmt
            selections.append(sel)
            count += 1

        self.search_selections = selections
        self.setExtraSelections(selections)
        return count

    def select_match_at(self, pos: int, length: int):
        """定位到指定匹配"""
        cursor = self.textCursor()
        cursor.setPosition(pos)
        cursor.setPosition(pos + length, QTextCursor.KeepAnchor)
        self.setTextCursor(cursor)
        self.ensureCursorVisible()
        self.setFocus()

    # ==================== 文档视图 ====================


class DocumentView(QWidget):
    """单个文档视图组件"""

    # 信号
    document_loaded = pyqtSignal(str, str, int)  # file_path, format, total_pages
    page_changed = pyqtSignal(int, int)  # current_page, total_pages
    zoom_changed = pyqtSignal(int)  # zoom_level

    SUPPORTED_FORMATS = {
        'pdf': 'PDF文档',
        'epub': 'EPUB电子书',
        'mobi': 'MOBI电子书',
        'fb2': 'FictionBook',
        'xps': 'XPS文档',
        'cbz': '漫画书',
        'docx': 'Word文档',
        'odt': 'OpenDocument文本',
        'odp': 'OpenDocument幻灯',
        'ods': 'OpenDocument表格',
        'txt': '文本文件',
        'md': 'Markdown',
        'pptx': 'PowerPoint',
    }

    def __init__(self, main_window):
        super().__init__(parent=main_window)
        self.main = main_window

        # 文档状态
        self.doc: Optional[fitz.Document] = None
        self.current_file: Optional[str] = None
        self.current_format: Optional[str] = None
        self.zoom_level = 100
        self.rotation = 0
        self.page_widgets: List[QWidget] = []
        self.render_threads: List[DocumentRenderThread] = []
        self.text_content: List[str] = []
        self.is_text_editable = False

        # 搜索状态
        self.search_term = ""
        self.search_results_pdf: Dict[int, List] = defaultdict(list)
        self.search_results_text: Dict[int, List[int]] = defaultdict(list)
        self.search_hits_linear: List[Tuple] = []
        self.current_search_index = -1
        self.current_search_hit: Optional[Tuple[int, int]] = None

        # 注释
        self.annotations: List[Dict] = []
        self.annotations_by_page: Dict[int, List[Dict]] = defaultdict(list)
        self.annotation_mode: Optional[str] = None

        # 书签
        self.bookmarks: List[Dict] = []

        # 笔记
        self.notes: List[Dict] = []

        # 阅读统计
        self.opened_at = time.time()
        self.total_read_seconds = 0.0
        self.current_page_index = 0

        # 历史记录（用于前进后退）
        self.page_history: List[int] = []
        self.history_index = -1

        # 自动保存定时器
        self.auto_save_timer = QTimer()
        self.auto_save_timer.timeout.connect(self._auto_save)

        # 构建UI
        self._build_ui()

    def _build_ui(self):
        """构建界面"""
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)

        # 主分割器
        splitter = QSplitter(Qt.Horizontal)
        layout.addWidget(splitter)

        # 左侧边栏
        self._build_sidebar()
        splitter.addWidget(self.sidebar)

        # 主阅读区
        self._build_reading_area()
        splitter.addWidget(self.scroll_area)

        splitter.setStretchFactor(0, 0)
        splitter.setStretchFactor(1, 1)
        splitter.setSizes([300, 1300])

    def _build_sidebar(self):
        """构建侧边栏"""
        self.sidebar = QTabWidget()
        self.sidebar.setMaximumWidth(350)
        self.sidebar.setMinimumWidth(250)

        # 目录
        self.outline_tree = QTreeWidget()
        self.outline_tree.setHeaderLabel("📑 目录")
        self.outline_tree.itemClicked.connect(self._on_outline_clicked)
        self.sidebar.addTab(self.outline_tree, "目录")

        # 缩略图
        self.thumbnail_list = QListWidget()
        self.thumbnail_list.setViewMode(QListWidget.IconMode)
        self.thumbnail_list.setIconSize(QSize(140, 180))
        self.thumbnail_list.setResizeMode(QListWidget.Adjust)
        self.thumbnail_list.itemClicked.connect(self._on_thumbnail_clicked)
        self.sidebar.addTab(self.thumbnail_list, "缩略图")

        # 书签
        self.bookmark_list = QListWidget()
        self.bookmark_list.setContextMenuPolicy(Qt.CustomContextMenu)
        self.bookmark_list.customContextMenuRequested.connect(self._show_bookmark_context_menu)
        self.bookmark_list.itemDoubleClicked.connect(self._jump_to_bookmark)
        self.sidebar.addTab(self.bookmark_list, "书签")

        # 注释
        self.annotation_list = QListWidget()
        self.annotation_list.setContextMenuPolicy(Qt.CustomContextMenu)
        self.annotation_list.customContextMenuRequested.connect(self._show_annotation_context_menu)
        self.annotation_list.itemDoubleClicked.connect(self._jump_to_annotation)
        self.sidebar.addTab(self.annotation_list, "注释")

        # 笔记
        self.notes_list = QListWidget()
        self.notes_list.setContextMenuPolicy(Qt.CustomContextMenu)
        self.notes_list.customContextMenuRequested.connect(self._show_note_context_menu)
        self.notes_list.itemDoubleClicked.connect(self._edit_note)
        self.sidebar.addTab(self.notes_list, "笔记")

    def _build_reading_area(self):
        """构建阅读区"""
        self.scroll_area = QScrollArea()
        self.scroll_area.setWidgetResizable(True)
        self.scroll_area.setHorizontalScrollBarPolicy(Qt.ScrollBarAsNeeded)
        self.scroll_area.setVerticalScrollBarPolicy(Qt.ScrollBarAlwaysOn)

        self.pages_container = QWidget()
        self.pages_layout = QVBoxLayout(self.pages_container)
        self.pages_layout.setAlignment(Qt.AlignTop | Qt.AlignHCenter)
        self.pages_layout.setSpacing(10)

        self.scroll_area.setWidget(self.pages_container)

        # 滚动监听
        self.scroll_area.verticalScrollBar().valueChanged.connect(self._on_scroll)

        # 平滑滚动支持
        self.scroll_animation: Optional[QPropertyAnimation] = None

        # ==================== 文档加载 ====================

    def load_document(self, file_path: Path, password: str = None):
        """加载文档"""
        try:
            # 清理之前的文档
            self._cleanup_current_document()

            file_path = Path(file_path)
            if not file_path.exists():
                raise FileNotFoundError(f"文件不存在: {file_path}")

            ext = file_path.suffix.lower()[1:]
            self.current_file = str(file_path)
            self.current_format = ext

            # 根据格式加载
            if ext in ['pdf', 'epub', 'mobi', 'azw', 'azw3', 'fb2', 'xps', 'cbz', 'cbr']:
                self._load_pymupdf_document(file_path, password)
            elif ext == 'docx' and DOCX_AVAILABLE:
                self._load_docx_document(file_path)
            elif ext == 'odt' and ODT_AVAILABLE:
                self._load_odt_document(file_path)
            elif ext == 'odp' and ODT_AVAILABLE:
                self._load_odp_document(file_path)
            elif ext == 'ods' and ODT_AVAILABLE:
                self._load_ods_document(file_path)
            elif ext == 'txt':
                self._load_text_document(file_path)
            elif ext in ['md', 'markdown']:
                self._load_markdown_document(file_path)
            elif ext == 'pptx' and PPTX_AVAILABLE:
                self._load_pptx_document(file_path)
            else:
                raise ValueError(f"不支持的文件格式: {ext}")

                # 加载文档状态
            self._load_document_state()

            # 发送信号
            total_pages = len(self.page_widgets)
            self.document_loaded.emit(self.current_file, ext, total_pages)

            # 启动自动保存
            config = ConfigManager.load()
            if config.get("auto_save", True):
                interval = config.get("auto_save_interval", 300) * 1000
                self.auto_save_timer.start(interval)

        except Exception as e:
            QMessageBox.critical(self, "加载失败", f"无法打开文件:\n{str(e)}")
            raise

    def _load_pymupdf_document(self, file_path: Path, password: str = None):
        """加载PyMuPDF支持的文档"""
        try:
            self.doc = fitz.open(str(file_path))

            # 检查是否需要密码
            if self.doc.needs_pass:
                if password is None:
                    password, ok = QInputDialog.getText(
                        self, "文档已加密", "请输入密码:",
                        QLineEdit.Password
                    )
                    if not ok or not password:
                        raise ValueError("文档已加密，需要密码")

                if not self.doc.authenticate(password):
                    raise ValueError("密码错误")

                    # 创建页面组件
            for i in range(len(self.doc)):
                page_widget = PageWidget(self, i)
                page_widget.annotation_created.connect(self._on_annotation_created)
                self.pages_layout.addWidget(page_widget)
                self.page_widgets.append(page_widget)

                # 加载目录和缩略图
            self._load_outline()
            self._load_thumbnails()

            # 渲染可见页面
            QTimer.singleShot(100, self._render_visible_pages)

        except Exception as e:
            raise Exception(f"PyMuPDF加载失败: {str(e)}")

    def _load_docx_document(self, file_path: Path):
        """加载Word文档"""
        doc = Document(str(file_path))
        full_text = []

        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)

        page_text = "\n\n".join(full_text)

        text_widget = TextPageWidget(self, 0, editable=True)
        text_widget.setPlainText(page_text)
        self.pages_layout.addWidget(text_widget)
        self.page_widgets.append(text_widget)
        self.text_content = [page_text]
        self.is_text_editable = True

    def _load_odt_document(self, file_path: Path):
        """加载ODT文档"""
        doc = odf_load(str(file_path))
        all_paras = doc.getElementsByType(odf_text.P)

        full_text = []
        for para in all_paras:
            para_text = teletype.extractText(para)
            if para_text.strip():
                full_text.append(para_text)

        page_text = "\n\n".join(full_text)

        text_widget = TextPageWidget(self, 0, editable=False)
        text_widget.setPlainText(page_text)
        self.pages_layout.addWidget(text_widget)
        self.page_widgets.append(text_widget)
        self.text_content = [page_text]

    def _load_odp_document(self, file_path: Path):
        """加载ODP文档"""
        self._load_odt_document(file_path)  # 简化处理

    def _load_ods_document(self, file_path: Path):
        """加载ODS文档"""
        self._load_odt_document(file_path)  # 简化处理

    def _load_text_document(self, file_path: Path):
        """加载文本文件"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()

        text_widget = TextPageWidget(self, 0, editable=True)
        text_widget.setPlainText(content)
        self.pages_layout.addWidget(text_widget)
        self.page_widgets.append(text_widget)
        self.text_content = [content]
        self.is_text_editable = True

    def _load_markdown_document(self, file_path: Path):
        """加载Markdown文档"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()

        text_widget = TextPageWidget(self, 0, editable=True)
        text_widget.setPlainText(content)
        self.pages_layout.addWidget(text_widget)
        self.page_widgets.append(text_widget)
        self.text_content = [content]
        self.is_text_editable = True

    def _load_pptx_document(self, file_path: Path):
        """加载PowerPoint文档"""
        prs = Presentation(str(file_path))

        for slide_num, slide in enumerate(prs.slides):
            slide_text = [f"=== 幻灯片 {slide_num + 1} ===\n"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            page_text = "\n\n".join(slide_text)

            text_widget = TextPageWidget(self, slide_num, editable=False)
            text_widget.setPlainText(page_text)
            self.pages_layout.addWidget(text_widget)
            self.page_widgets.append(text_widget)
            self.text_content.append(page_text)

    def _cleanup_current_document(self):
        """清理当前文档"""
        # 停止所有渲染线程
        for thread in self.render_threads:
            thread.stop()
            thread.wait()
        self.render_threads.clear()

        # 清理页面组件
        for widget in self.page_widgets:
            widget.deleteLater()
        self.page_widgets.clear()

        # 清理数据
        self.text_content.clear()
        self.search_results_pdf.clear()
        self.search_results_text.clear()
        self.search_hits_linear.clear()
        self.annotations.clear()
        self.annotations_by_page.clear()
        self.bookmarks.clear()
        self.notes.clear()
        self.page_history.clear()

        # 清理UI
        self.outline_tree.clear()
        self.thumbnail_list.clear()
        self.bookmark_list.clear()
        self.annotation_list.clear()
        self.notes_list.clear()

        # 关闭文档
        if self.doc:
            try:
                self.doc.close()
            except:
                pass
            self.doc = None

            # 停止自动保存
        self.auto_save_timer.stop()

        # ==================== 渲染与显示 ====================

    def _render_visible_pages(self):
        """渲染可见页面"""
        if not self.page_widgets or not isinstance(self.page_widgets[0], PageWidget):
            return

        if not self.doc:
            return

        viewport = self.scroll_area.viewport().rect()
        scroll_value = self.scroll_area.verticalScrollBar().value()

        config = ConfigManager.load()
        preload = config.get("preload_pages", 3)

        for i, widget in enumerate(self.page_widgets):
            if not isinstance(widget, PageWidget):
                continue

            widget_pos = widget.pos().y() - scroll_value
            widget_height = widget.height()

            # 判断是否在可见范围或预加载范围
            in_viewport = -widget_height <= widget_pos <= viewport.height() + widget_height
            in_preload = abs(i - self.current_page_index) <= preload

            if (in_viewport or in_preload) and (widget.pixmap() is None or widget.pixmap().isNull()):
                self._render_page(i)

    def _render_page(self, page_num: int):
        """渲染指定页面"""
        if not self.doc or page_num >= len(self.doc):
            return

        theme = self.main.current_theme
        dark_mode = (theme == "dark")

        thread = DocumentRenderThread(
            self.doc, page_num, self.zoom_level,
            self.rotation, dark_mode, high_quality=True
        )
        thread.page_rendered.connect(self._on_page_rendered)
        thread.render_failed.connect(self._on_render_failed)
        thread.start()
        self.render_threads.append(thread)

    @pyqtSlot(int, QPixmap)
    def _on_page_rendered(self, page_num: int, pixmap: QPixmap):
        """页面渲染完成"""
        if page_num < len(self.page_widgets):
            widget = self.page_widgets[page_num]
            if isinstance(widget, PageWidget):
                widget.setPixmap(pixmap)
                widget.setFixedSize(pixmap.size())

    @pyqtSlot(int, str)
    def _on_render_failed(self, page_num: int, error: str):
        """渲染失败"""
        print(f"页面 {page_num + 1} 渲染失败: {error}")

    def _on_scroll(self):
        """滚动事件处理"""
        self._render_visible_pages()
        self._update_page_info()

        # 延迟保存阅读位置
        if hasattr(self, '_save_timer'):
            self._save_timer.stop()

        self._save_timer = QTimer()
        self._save_timer.setSingleShot(True)
        self._save_timer.timeout.connect(self._save_reading_position)
        self._save_timer.start(1000)

    def _update_page_info(self):
        """更新当前页码"""
        if not self.page_widgets:
            return

        viewport_center = self.scroll_area.viewport().height() // 2
        scroll_value = self.scroll_area.verticalScrollBar().value()

        for i, widget in enumerate(self.page_widgets):
            widget_top = widget.pos().y() - scroll_value
            widget_bottom = widget_top + widget.height()

            if widget_top <= viewport_center <= widget_bottom:
                if self.current_page_index != i:
                    self.current_page_index = i
                    self.page_changed.emit(i + 1, len(self.page_widgets))

                    # 添加到历史
                    if not self.page_history or self.page_history[-1] != i:
                        self.page_history.append(i)
                        self.history_index = len(self.page_history) - 1
                break

                # ==================== 视图控制 ====================

    def zoom_in(self):
        """放大"""
        self.zoom_level = min(self.zoom_level + 25, 500)
        self._update_zoom()

    def zoom_out(self):
        """缩小"""
        self.zoom_level = max(self.zoom_level - 25, 25)
        self._update_zoom()

    def set_zoom(self, zoom: int):
        """设置缩放级别"""
        self.zoom_level = max(25, min(500, zoom))
        self._update_zoom()

    def fit_width(self):
        """适合宽度"""
        if not self.doc or not self.page_widgets:
            return

        try:
            page = self.doc[0]
            page_width = page.rect.width
            viewport_width = self.scroll_area.viewport().width() - 40
            self.zoom_level = int((viewport_width / page_width) * 100)
            self._update_zoom()
        except:
            pass

    def fit_page(self):
        """适合页面"""
        if not self.doc or not self.page_widgets:
            return

        try:
            page = self.doc[0]
            page_width = page.rect.width
            page_height = page.rect.height
            viewport_width = self.scroll_area.viewport().width() - 40
            viewport_height = self.scroll_area.viewport().height() - 40

            zoom_w = (viewport_width / page_width) * 100
            zoom_h = (viewport_height / page_height) * 100
            self.zoom_level = int(min(zoom_w, zoom_h))
            self._update_zoom()
        except:
            pass

    def _update_zoom(self):
        """更新缩放"""
        self.zoom_changed.emit(self.zoom_level)

        if self.page_widgets and isinstance(self.page_widgets[0], TextPageWidget):
            # 文本文档缩放
            base_size = 11
            new_size = int(base_size * self.zoom_level / 100)
            for widget in self.page_widgets:
                if isinstance(widget, TextPageWidget):
                    font = widget.font()
                    font.setPointSize(max(6, min(72, new_size)))
                    widget.setFont(font)
        else:
            # 图像文档重新渲染
            for widget in self.page_widgets:
                if isinstance(widget, PageWidget):
                    widget.clear()
            QTimer.singleShot(50, self._render_visible_pages)

    def rotate_page(self):
        """旋转页面"""
        self.rotation = (self.rotation + 90) % 360

        for widget in self.page_widgets:
            if isinstance(widget, PageWidget):
                widget.clear()

        QTimer.singleShot(50, self._render_visible_pages)

    def toggle_theme(self, theme: str):
        """切换主题"""
        for widget in self.page_widgets:
            if isinstance(widget, TextPageWidget):
                widget._apply_theme()
            elif isinstance(widget, PageWidget):
                widget.clear()

        if isinstance(self.page_widgets[0], PageWidget) if self.page_widgets else False:
            QTimer.singleShot(50, self._render_visible_pages)

            # ==================== 页面导航 ====================

    def previous_page(self):
        """上一页"""
        if not self.page_widgets:
            return
        new_index = max(0, self.current_page_index - 1)
        self._scroll_to_page(new_index)

    def next_page(self):
        """下一页"""
        if not self.page_widgets:
            return
        new_index = min(len(self.page_widgets) - 1, self.current_page_index + 1)
        self._scroll_to_page(new_index)

    def jump_to_page(self, page_index: int):
        """跳转到指定页"""
        if 0 <= page_index < len(self.page_widgets):
            self._scroll_to_page(page_index, smooth=False)

    def go_back(self):
        """后退"""
        if self.history_index > 0:
            self.history_index -= 1
            page = self.page_history[self.history_index]
            self._scroll_to_page(page, add_to_history=False)

    def go_forward(self):
        """前进"""
        if self.history_index < len(self.page_history) - 1:
            self.history_index += 1
            page = self.page_history[self.history_index]
            self._scroll_to_page(page, add_to_history=False)

    def _scroll_to_page(self, page_num: int, smooth: bool = True, add_to_history: bool = True):
        """滚动到指定页面"""
        if not (0 <= page_num < len(self.page_widgets)):
            return

        widget = self.page_widgets[page_num]
        target_value = widget.pos().y()

        config = ConfigManager.load()
        use_smooth = smooth and config.get("smooth_scroll", True)

        if use_smooth:
            # 平滑滚动
            if self.scroll_animation:
                self.scroll_animation.stop()

            scrollbar = self.scroll_area.verticalScrollBar()
            self.scroll_animation = QPropertyAnimation(scrollbar, b"value")
            self.scroll_animation.setDuration(300)
            self.scroll_animation.setStartValue(scrollbar.value())
            self.scroll_animation.setEndValue(target_value)
            self.scroll_animation.setEasingCurve(QEasingCurve.OutCubic)
            self.scroll_animation.start()
        else:
            # 直接跳转
            self.scroll_area.verticalScrollBar().setValue(target_value)

        if add_to_history:
            if not self.page_history or self.page_history[-1] != page_num:
                # 清除前进历史
                self.page_history = self.page_history[:self.history_index + 1]
                self.page_history.append(page_num)
                self.history_index = len(self.page_history) - 1

        self.current_page_index = page_num
        self.page_changed.emit(page_num + 1, len(self.page_widgets))

        # ==================== 目录和缩略图 ====================

    def _load_outline(self):
        """加载目录"""
        self.outline_tree.clear()
        if not self.doc:
            return

        try:
            toc = self.doc.get_toc()
            if not toc:
                return

            items_stack = [self.outline_tree.invisibleRootItem()]

            for level, title, page in toc:
                item = QTreeWidgetItem([f"{title} (第{page}页)"])
                item.setData(0, Qt.UserRole, page - 1)

                while len(items_stack) > level:
                    items_stack.pop()

                items_stack[-1].addChild(item)
                items_stack.append(item)

            self.outline_tree.expandAll()
        except:
            pass

    def _on_outline_clicked(self, item: QTreeWidgetItem):
        """目录项点击"""
        page_num = item.data(0, Qt.UserRole)
        if page_num is not None:
            self._scroll_to_page(page_num)

    def _load_thumbnails(self):
        """加载缩略图"""
        self.thumbnail_list.clear()
        if not self.doc:
            return

        try:
            # 限制缩略图数量
            max_thumbs = min(len(self.doc), 100)

            for i in range(max_thumbs):
                page = self.doc[i]
                mat = fitz.Matrix(0.2, 0.2)
                pix = page.get_pixmap(matrix=mat, alpha=False)

                img = QImage(
                    pix.samples, pix.width, pix.height,
                    pix.stride, QImage.Format_RGB888
                )
                pixmap = QPixmap.fromImage(img)

                item = QListWidgetItem(QIcon(pixmap), f"第{i + 1}页")
                item.setData(Qt.UserRole, i)
                self.thumbnail_list.addItem(item)
        except:
            pass

    def _on_thumbnail_clicked(self, item: QListWidgetItem):
        """缩略图点击"""
        page_num = item.data(Qt.UserRole)
        self._scroll_to_page(page_num)

        # ==================== 书签 ====================

    def add_bookmark(self, title: str = None):
        """添加书签"""
        if not self.page_widgets:
            return

        page = self.current_page_index

        if title is None:
            title, ok = QInputDialog.getText(
                self, "添加书签",
                f"为第 {page + 1} 页添加书签:",
                text=f"第 {page + 1} 页"
            )
            if not ok or not title.strip():
                return

        bookmark = {
            "page": page,
            "title": title.strip(),
            "created": time.time()
        }

        self.bookmarks.append(bookmark)

        item = QListWidgetItem(f"⭐ {title} (第{page + 1}页)")
        item.setData(Qt.UserRole, bookmark)
        self.bookmark_list.addItem(item)

        self._save_document_state()
        self.main.show_status(f"已添加书签: {title}", 3000)

    def _jump_to_bookmark(self, item: QListWidgetItem):
        """跳转到书签"""
        bookmark = item.data(Qt.UserRole)
        if isinstance(bookmark, dict):
            self._scroll_to_page(bookmark.get("page", 0))

    def _show_bookmark_context_menu(self, pos: QPoint):
        """书签右键菜单"""
        item = self.bookmark_list.itemAt(pos)
        if not item:
            return

        menu = QMenu(self)

        rename_action = menu.addAction("✏️ 重命名")
        delete_action = menu.addAction("🗑️ 删除")

        action = menu.exec_(self.bookmark_list.mapToGlobal(pos))

        if action == rename_action:
            self._rename_bookmark(item)
        elif action == delete_action:
            self._delete_bookmark(item)

    def _rename_bookmark(self, item: QListWidgetItem):
        """重命名书签"""
        bookmark = item.data(Qt.UserRole)
        if not isinstance(bookmark, dict):
            return

        new_title, ok = QInputDialog.getText(
            self, "重命名书签", "新标题:",
            text=bookmark.get("title", "")
        )

        if ok and new_title.strip():
            bookmark["title"] = new_title.strip()
            page = bookmark.get("page", 0)
            item.setText(f"⭐ {new_title} (第{page + 1}页)")
            self._save_document_state()

    def _delete_bookmark(self, item: QListWidgetItem):
        """删除书签"""
        bookmark = item.data(Qt.UserRole)
        if isinstance(bookmark, dict) and bookmark in self.bookmarks:
            self.bookmarks.remove(bookmark)
            self.bookmark_list.takeItem(self.bookmark_list.row(item))
            self._save_document_state()

            # ==================== 注释 ====================

    def set_annotation_mode(self, mode: Optional[str]):
        """设置注释模式"""
        self.annotation_mode = mode

        if mode:
            self.main.show_status(f"注释模式: {mode}", 3000)
        else:
            self.main.show_status("已退出注释模式", 2000)

    @pyqtSlot(int, QRect, str)
    def _on_annotation_created(self, page_num: int, widget_rect: QRect, ann_type: str):
        """创建注释"""
        if not self.doc or self.current_format != 'pdf':
            return

            # 转换为PDF坐标
        zoom = self.zoom_level / 100.0
        x0 = widget_rect.left() / zoom
        y0 = widget_rect.top() / zoom
        x1 = widget_rect.right() / zoom
        y1 = widget_rect.bottom() / zoom

        content = ""
        if ann_type == "note":
            content, ok = QInputDialog.getMultiLineText(
                self, "添加批注", "批注内容:"
            )
            if not ok:
                return

        annotation = {
            "page": page_num,
            "type": ann_type,
            "rect": [float(x0), float(y0), float(x1), float(y1)],
            "content": content,
            "created": time.time()
        }

        self.annotations.append(annotation)
        self.annotations_by_page[page_num].append(annotation)
        self._add_annotation_to_list(annotation)
        self._save_document_state()

        # 刷新页面
        widget = self.page_widgets[page_num]
        if isinstance(widget, PageWidget):
            widget.update()

        self.main.show_status("注释已添加", 2000)

    def _add_annotation_to_list(self, annotation: Dict):
        """添加注释到列表"""
        page = annotation.get("page", 0) + 1
        ann_type = annotation.get("type", "highlight")
        content = annotation.get("content", "").strip()

        type_name = {
            "highlight": "🖍 高亮",
            "underline": "〰 下划线",
            "note": "💬 批注"
        }.get(ann_type, ann_type)

        if content:
            display = f"{type_name} - 第{page}页\n{content[:40]}..."
        else:
            display = f"{type_name} - 第{page}页"

        item = QListWidgetItem(display)
        item.setData(Qt.UserRole, annotation)
        self.annotation_list.addItem(item)

    def _jump_to_annotation(self, item: QListWidgetItem):
        """跳转到注释"""
        annotation = item.data(Qt.UserRole)
        if isinstance(annotation, dict):
            page = annotation.get("page", 0)
            self._scroll_to_page(page)

    def _show_annotation_context_menu(self, pos: QPoint):
        """注释右键菜单"""
        item = self.annotation_list.itemAt(pos)
        if not item:
            return

        menu = QMenu(self)

        edit_action = menu.addAction("✏️ 编辑")
        delete_action = menu.addAction("🗑️ 删除")

        action = menu.exec_(self.annotation_list.mapToGlobal(pos))

        if action == edit_action:
            self._edit_annotation(item)
        elif action == delete_action:
            self._delete_annotation(item)

    def _edit_annotation(self, item: QListWidgetItem):
        """编辑注释"""
        annotation = item.data(Qt.UserRole)
        if not isinstance(annotation, dict):
            return

        if annotation.get("type") == "note":
            new_content, ok = QInputDialog.getMultiLineText(
                self, "编辑批注", "批注内容:",
                annotation.get("content", "")
            )

            if ok:
                annotation["content"] = new_content
                self._save_document_state()

                # 更新显示
                page = annotation.get("page", 0) + 1
                display = f"💬 批注 - 第{page}页\n{new_content[:40]}..."
                item.setText(display)

    def _delete_annotation(self, item: QListWidgetItem):
        """删除注释"""
        annotation = item.data(Qt.UserRole)
        if not isinstance(annotation, dict):
            return

        page = annotation.get("page", 0)

        if annotation in self.annotations:
            self.annotations.remove(annotation)

        if annotation in self.annotations_by_page[page]:
            self.annotations_by_page[page].remove(annotation)

        self.annotation_list.takeItem(self.annotation_list.row(item))
        self._save_document_state()

        # 刷新页面
        if page < len(self.page_widgets):
            widget = self.page_widgets[page]
            if isinstance(widget, PageWidget):
                widget.update()

                # ==================== 笔记系统 ====================

    def add_note(self):
        """添加笔记"""
        content, ok = QInputDialog.getMultiLineText(
            self, "添加笔记", "笔记内容:"
        )

        if not ok or not content.strip():
            return

        note = {
            "page": self.current_page_index,
            "content": content.strip(),
            "created": time.time()
        }

        self.notes.append(note)

        page = note.get("page", 0) + 1
        preview = content.strip().split('\n')[0][:30]
        item = QListWidgetItem(f"📝 第{page}页 - {preview}...")
        item.setData(Qt.UserRole, note)
        self.notes_list.addItem(item)

        self._save_document_state()
        self.main.show_status("笔记已添加", 2000)

    def _edit_note(self, item: QListWidgetItem):
        """编辑笔记"""
        note = item.data(Qt.UserRole)
        if not isinstance(note, dict):
            return

        content, ok = QInputDialog.getMultiLineText(
            self, "编辑笔记", "笔记内容:",
            note.get("content", "")
        )

        if ok:
            note["content"] = content.strip()
            self._save_document_state()

            page = note.get("page", 0) + 1
            preview = content.strip().split('\n')[0][:30]
            item.setText(f"📝 第{page}页 - {preview}...")

    def _show_note_context_menu(self, pos: QPoint):
        """笔记右键菜单"""
        item = self.notes_list.itemAt(pos)
        if not item:
            return

        menu = QMenu(self)

        edit_action = menu.addAction("✏️ 编辑")
        delete_action = menu.addAction("🗑️ 删除")
        jump_action = menu.addAction("🔍 定位")

        action = menu.exec_(self.notes_list.mapToGlobal(pos))

        if action == edit_action:
            self._edit_note(item)
        elif action == delete_action:
            self._delete_note(item)
        elif action == jump_action:
            note = item.data(Qt.UserRole)
            if isinstance(note, dict):
                self._scroll_to_page(note.get("page", 0))

    def _delete_note(self, item: QListWidgetItem):
        """删除笔记"""
        note = item.data(Qt.UserRole)
        if isinstance(note, dict) and note in self.notes:
            self.notes.remove(note)
            self.notes_list.takeItem(self.notes_list.row(item))
            self._save_document_state()

            # ==================== 搜索 ====================

    def search_text(self, term: str):
        """搜索文本"""
        self.search_term = term
        self.search_results_pdf.clear()
        self.search_results_text.clear()
        self.search_hits_linear.clear()
        self.current_search_hit = None
        self.current_search_index = -1

        if not term:
            # 清除高亮
            for widget in self.page_widgets:
                if isinstance(widget, TextPageWidget):
                    widget.clear_search_highlight()
                elif isinstance(widget, PageWidget):
                    widget.update()
            return

        total_hits = 0

        # 文本文档搜索
        if self.text_content:
            for i, widget in enumerate(self.page_widgets):
                if isinstance(widget, TextPageWidget):
                    count = widget.highlight_all(term, QColor(255, 255, 0, 150))

                    if count > 0:
                        doc = widget.document()
                        cursor = QTextCursor(doc)
                        positions = []

                        while True:
                            cursor = doc.find(term, cursor)
                            if cursor.isNull():
                                break
                            positions.append(cursor.selectionStart())

                        self.search_results_text[i] = positions

                        for idx in range(len(positions)):
                            self.search_hits_linear.append(("text", i, idx))

                        total_hits += count

                        # PDF文档搜索
        elif self.doc:
            for page_num in range(len(self.doc)):
                page = self.doc[page_num]
                rects = page.search_for(term)

                if rects:
                    self.search_results_pdf[page_num] = rects

                    for idx in range(len(rects)):
                        self.search_hits_linear.append(("pdf", page_num, idx))

                    total_hits += len(rects)

                    # 刷新页面以显示高亮
            for widget in self.page_widgets:
                if isinstance(widget, PageWidget):
                    widget.update()

        if total_hits == 0:
            self.main.show_status("未找到匹配结果", 3000)
            return

            # 跳转到第一个匹配
        self.current_search_index = -1
        self.search_next()
        self.main.show_status(f"找到 {total_hits} 个匹配", 5000)

    def search_next(self):
        """下一个搜索结果"""
        if not self.search_hits_linear:
            return

        self.current_search_index = (self.current_search_index + 1) % len(self.search_hits_linear)
        self._goto_current_search_hit()

    def search_prev(self):
        """上一个搜索结果"""
        if not self.search_hits_linear:
            return

        self.current_search_index = (self.current_search_index - 1) % len(self.search_hits_linear)
        self._goto_current_search_hit()

    def _goto_current_search_hit(self):
        """跳转到当前搜索结果"""
        if not (0 <= self.current_search_index < len(self.search_hits_linear)):
            return

        kind, page, idx = self.search_hits_linear[self.current_search_index]

        if kind == "text":
            widget = self.page_widgets[page]
            positions = self.search_results_text.get(page, [])

            if positions and idx < len(positions):
                pos = positions[idx]
                self._scroll_to_page(page)
                widget.select_match_at(pos, len(self.search_term))

        elif kind == "pdf":
            self.current_search_hit = (page, idx)
            self._scroll_to_page(page)

            # 刷新页面以突出当前匹配
            widget = self.page_widgets[page]
            if isinstance(widget, PageWidget):
                widget.update()

                # ==================== 导出 ====================

    def export_pages(self):
        """导出页面"""
        if not self.doc:
            QMessageBox.information(
                self, "导出页面",
                "当前文档类型不支持导出功能"
            )
            return

            # 页码输入
        text, ok = QInputDialog.getText(
            self, "导出页面",
            "输入页码范围（例如: 1-5 或 1,3,5）:",
            text=f"1-{len(self.doc)}"
        )

        if not ok or not text.strip():
            return

            # 解析页码
        pages = self._parse_page_range(text)

        if not pages:
            QMessageBox.warning(self, "错误", "页码格式不正确或没有有效页码")
            return

            # 选择导出类型
        export_type, ok = QInputDialog.getItem(
            self, "导出类型", "选择导出格式:",
            ["PDF文档", "PNG图片", "JPEG图片"], 0, False
        )

        if not ok:
            return

        if export_type == "PDF文档":
            self._export_as_pdf(pages)
        else:
            fmt = "PNG" if export_type == "PNG图片" else "JPEG"
            self._export_as_images(pages, fmt)

    def _parse_page_range(self, text: str) -> List[int]:
        """解析页码范围"""
        pages = set()
        parts = text.replace("，", ",").split(",")

        try:
            for part in parts:
                part = part.strip()

                if "-" in part:
                    start, end = part.split("-")
                    start = int(start) - 1
                    end = int(end) - 1

                    for p in range(start, end + 1):
                        if 0 <= p < len(self.doc):
                            pages.add(p)
                else:
                    p = int(part) - 1
                    if 0 <= p < len(self.doc):
                        pages.add(p)
        except:
            return []

        return sorted(pages)

    def _export_as_pdf(self, pages: List[int]):
        """导出为PDF"""
        save_path, _ = QFileDialog.getSaveFileName(
            self, "导出PDF", "", "PDF文件 (*.pdf)"
        )

        if not save_path:
            return

        try:
            new_doc = fitz.open()

            for page_num in pages:
                new_doc.insert_pdf(self.doc, from_page=page_num, to_page=page_num)

            new_doc.save(save_path)
            new_doc.close()

            self.main.show_status("PDF导出成功", 4000)
            QMessageBox.information(self, "导出成功", f"已导出到:\n{save_path}")

        except Exception as e:
            QMessageBox.critical(self, "导出失败", str(e))

    def _export_as_images(self, pages: List[int], fmt: str):
        """导出为图片"""
        dir_path = QFileDialog.getExistingDirectory(self, "选择导出目录")

        if not dir_path:
            return

        try:
            progress = QProgressDialog(
                "正在导出图片...", "取消", 0, len(pages), self
            )
            progress.setWindowModality(Qt.WindowModal)

            for i, page_num in enumerate(pages):
                if progress.wasCanceled():
                    break

                page = self.doc[page_num]
                mat = fitz.Matrix(3, 3)  # 高分辨率
                pix = page.get_pixmap(matrix=mat, alpha=False)

                ext = "png" if fmt == "PNG" else "jpg"
                output = Path(dir_path) / f"page_{page_num + 1}.{ext}"
                pix.save(str(output))

                progress.setValue(i + 1)

            progress.close()
            self.main.show_status(f"图片导出成功（{len(pages)}页）", 4000)
            QMessageBox.information(self, "导出成功", f"已导出到:\n{dir_path}")

        except Exception as e:
            QMessageBox.critical(self, "导出失败", str(e))

            # ==================== OCR识别 ====================

    def ocr_current_page(self):
        """OCR识别当前页"""
        if not OCR_AVAILABLE:
            QMessageBox.information(
                self, "OCR不可用",
                "请安装 tesseract-ocr 和 pytesseract 库以使用OCR功能"
            )
            return

        if not self.doc:
            return

        try:
            page_num = self.current_page_index
            page = self.doc[page_num]

            # 渲染为图片
            mat = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat, alpha=False)

            # 转换为PIL Image
            img_data = pix.tobytes("png")
            pil_image = Image.open(io.BytesIO(img_data))

            # OCR识别
            progress = QProgressDialog("正在识别文字...", None, 0, 0, self)
            progress.setWindowModality(Qt.WindowModal)
            progress.show()
            QApplication.processEvents()

            config = ConfigManager.load()
            lang = config.get("ocr_language", "chi_sim+eng")

            text = pytesseract.image_to_string(pil_image, lang=lang)

            progress.close()

            # 显示结果
            self._show_ocr_result(text, page_num)

        except Exception as e:
            QMessageBox.critical(self, "OCR失败", str(e))

    def _show_ocr_result(self, text: str, page_num: int):
        """显示OCR结果"""
        dialog = QDialog(self)
        dialog.setWindowTitle(f"OCR识别结果 - 第{page_num + 1}页")
        dialog.setMinimumSize(600, 400)

        layout = QVBoxLayout(dialog)

        text_edit = QTextEdit()
        text_edit.setPlainText(text)
        layout.addWidget(text_edit)

        button_layout = QHBoxLayout()

        copy_btn = QPushButton("📋 复制")
        copy_btn.clicked.connect(lambda: QApplication.clipboard().setText(text))
        button_layout.addWidget(copy_btn)

        save_btn = QPushButton("💾 保存")
        save_btn.clicked.connect(lambda: self._save_ocr_result(text, page_num))
        button_layout.addWidget(save_btn)

        close_btn = QPushButton("关闭")
        close_btn.clicked.connect(dialog.close)
        button_layout.addWidget(close_btn)

        layout.addLayout(button_layout)

        dialog.exec_()

    def _save_ocr_result(self, text: str, page_num: int):
        """保存OCR结果"""
        save_path, _ = QFileDialog.getSaveFileName(
            self, "保存OCR结果",
            f"page_{page_num + 1}_ocr.txt",
            "文本文件 (*.txt)"
        )

        if save_path:
            try:
                with open(save_path, 'w', encoding='utf-8') as f:
                    f.write(text)
                self.main.show_status("OCR结果已保存", 3000)
            except Exception as e:
                QMessageBox.critical(self, "保存失败", str(e))

                # ==================== 打印 ====================

    def print_document(self):
        """打印文档"""
        if not self.doc and not self.text_content:
            return

        printer = QPrinter(QPrinter.HighResolution)
        dialog = QPrintDialog(printer, self)

        if dialog.exec_() == QPrintDialog.Accepted:
            self._do_print(printer)

    def print_preview(self):
        """打印预览"""
        if not self.doc and not self.text_content:
            return

        printer = QPrinter(QPrinter.HighResolution)
        preview = QPrintPreviewDialog(printer, self)
        preview.paintRequested.connect(self._do_print)
        preview.exec_()

    def _do_print(self, printer: QPrinter):
        """执行打印"""
        painter = QPainter()
        painter.begin(printer)

        try:
            if self.doc:
                # PDF文档打印
                for i in range(len(self.doc)):
                    if i > 0:
                        printer.newPage()

                    page = self.doc[i]
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))

                    img = QImage(
                        pix.samples, pix.width, pix.height,
                        pix.stride, QImage.Format_RGB888
                    )

                    pixmap = QPixmap.fromImage(img)

                    # 缩放到页面大小
                    rect = printer.pageRect()
                    scaled = pixmap.scaled(
                        rect.width(), rect.height(),
                        Qt.KeepAspectRatio, Qt.SmoothTransformation
                    )

                    painter.drawPixmap(0, 0, scaled)

            elif self.text_content:
                # 文本文档打印
                document = QTextDocument()
                document.setPlainText("\n\n".join(self.text_content))
                document.print_(printer)

        finally:
            painter.end()

            # ==================== 状态持久化 ====================

    def _load_document_state(self):
        """加载文档状态"""
        if not self.current_file:
            return

        state = ConfigManager.get_document_state(self.current_file)

        # 书签
        self.bookmarks = state.get("bookmarks", [])
        for bm in self.bookmarks:
            title = bm.get("title", f"第{bm['page'] + 1}页")
            page = bm.get("page", 0)
            item = QListWidgetItem(f"⭐ {title} (第{page + 1}页)")
            item.setData(Qt.UserRole, bm)
            self.bookmark_list.addItem(item)

            # 注释
        self.annotations = state.get("annotations", [])
        self.annotations_by_page.clear()
        for ann in self.annotations:
            page = ann.get("page", 0)
            self.annotations_by_page[page].append(ann)
            self._add_annotation_to_list(ann)

            # 笔记
        self.notes = state.get("notes", [])
        for note in self.notes:
            page = note.get("page", 0) + 1
            content = note.get("content", "")
            preview = content.split('\n')[0][:30]
            item = QListWidgetItem(f"📝 第{page}页 - {preview}...")
            item.setData(Qt.UserRole, note)
            self.notes_list.addItem(item)

            # 阅读统计
        stats = state.get("stats", {})
        self.total_read_seconds = float(stats.get("total_seconds", 0.0))

        # 恢复阅读位置
        QTimer.singleShot(500, self._restore_reading_position)

    def _save_document_state(self):
        """保存文档状态"""
        if not self.current_file:
            return

        elapsed = time.time() - self.opened_at
        total = self.total_read_seconds + elapsed

        state = {
            "bookmarks": self.bookmarks,
            "annotations": self.annotations,
            "notes": self.notes,
            "stats": {
                "total_seconds": total,
                "last_page": self.current_page_index
            }
        }

        ConfigManager.set_document_state(self.current_file, state)

    def _save_reading_position(self):
        """保存阅读位置"""
        if not self.current_file:
            return

        config = ConfigManager.load()
        positions = config.get("reading_positions", {})

        positions[self.current_file] = {
            "page": self.current_page_index,
            "scroll": self.scroll_area.verticalScrollBar().value(),
            "zoom": self.zoom_level
        }

        config["reading_positions"] = positions
        ConfigManager.save(config)

    def _restore_reading_position(self):
        """恢复阅读位置"""
        if not self.current_file:
            return

        config = ConfigManager.load()
        positions = config.get("reading_positions", {})
        pos = positions.get(self.current_file)

        if not pos:
            return

            # 恢复缩放
        zoom = pos.get("zoom", 100)
        if zoom != self.zoom_level:
            self.zoom_level = zoom
            self._update_zoom()

            # 恢复页面位置
        page = pos.get("page", 0)
        self._scroll_to_page(page, smooth=False)

    def _auto_save(self):
        """自动保存"""
        self._save_document_state()
        self._save_reading_position()

        # ==================== 清理 ====================

    def close_view(self):
        """关闭视图"""
        self._save_document_state()
        self._save_reading_position()
        self._cleanup_current_document()

    def get_progress_info(self) -> Dict:
        """获取阅读进度信息"""
        total_pages = len(self.page_widgets)
        current = self.current_page_index + 1

        elapsed = time.time() - self.opened_at
        total_time = self.total_read_seconds + elapsed

        hours = int(total_time // 3600)
        minutes = int((total_time % 3600) // 60)
        seconds = int(total_time % 60)

        return {
            "current_page": current,
            "total_pages": total_pages,
            "progress_percent": (current / total_pages * 100) if total_pages > 0 else 0,
            "time_hours": hours,
            "time_minutes": minutes,
            "time_seconds": seconds,
            "time_formatted": f"{hours:02d}:{minutes:02d}:{seconds:02d}"
        }

    # ==================== 主窗口 ====================


class UniversalDocumentReader(QMainWindow):
    """主窗口"""

    def __init__(self):
        super().__init__()

        self.setWindowTitle("专业文档阅读器")
        self.setAcceptDrops(True)

        # 当前主题
        self.current_theme = "light"

        # 分离的窗口列表
        self.detached_windows: List['UniversalDocumentReader'] = []

        # 构建UI
        self._build_ui()

        # 加载配置
        self._load_config()

        # 设置快捷键
        self._setup_shortcuts()

        # 定时更新阅读信息
        self.update_timer = QTimer()
        self.update_timer.timeout.connect(self._update_reading_info)
        self.update_timer.start(1000)

    def _build_ui(self):
        """构建用户界面"""
        # 中央标签页
        self.tab_widget = QTabWidget()
        self.tab_widget.setTabsClosable(True)
        self.tab_widget.setMovable(True)
        self.tab_widget.setDocumentMode(True)
        self.tab_widget.currentChanged.connect(self._on_tab_changed)
        self.tab_widget.tabCloseRequested.connect(self._close_tab)
        self.tab_widget.tabBarDoubleClicked.connect(self._detach_tab)

        self.setCentralWidget(self.tab_widget)

        # 菜单栏
        self._build_menubar()

        # 工具栏
        self._build_toolbar()

        # 状态栏
        self._build_statusbar()

    def _build_menubar(self):
        """构建菜单栏"""
        menubar = self.menuBar()

        # 文件菜单
        file_menu = menubar.addMenu("📂 文件")

        open_action = QAction("打开文档...", self)
        open_action.setShortcut(QKeySequence.Open)
        open_action.triggered.connect(self.open_file)
        file_menu.addAction(open_action)

        recent_menu = file_menu.addMenu("🕐 最近打开")
        recent_menu.aboutToShow.connect(lambda: self._populate_recent_menu(recent_menu))

        file_menu.addSeparator()

        save_action = QAction("保存", self)
        save_action.setShortcut(QKeySequence.Save)
        save_action.triggered.connect(self.save_current_document)
        file_menu.addAction(save_action)

        export_action = QAction("导出页面...", self)
        export_action.triggered.connect(self.export_pages)
        file_menu.addAction(export_action)

        file_menu.addSeparator()

        print_preview_action = QAction("打印预览...", self)
        print_preview_action.triggered.connect(self.print_preview)
        file_menu.addAction(print_preview_action)

        print_action = QAction("打印...", self)
        print_action.setShortcut(QKeySequence.Print)
        print_action.triggered.connect(self.print_document)
        file_menu.addAction(print_action)

        file_menu.addSeparator()

        close_tab_action = QAction("关闭标签页", self)
        close_tab_action.setShortcut(QKeySequence.Close)
        close_tab_action.triggered.connect(lambda: self._close_tab(self.tab_widget.currentIndex()))
        file_menu.addAction(close_tab_action)

        quit_action = QAction("退出", self)
        quit_action.setShortcut(QKeySequence.Quit)
        quit_action.triggered.connect(self.close)
        file_menu.addAction(quit_action)

        # 编辑菜单
        edit_menu = menubar.addMenu("✏️ 编辑")

        find_action = QAction("查找...", self)
        find_action.setShortcut(QKeySequence.Find)
        find_action.triggered.connect(lambda: self.search_input.setFocus())
        edit_menu.addAction(find_action)

        find_next_action = QAction("查找下一个", self)
        find_next_action.setShortcut(QKeySequence.FindNext)
        find_next_action.triggered.connect(self.search_next)
        edit_menu.addAction(find_next_action)

        find_prev_action = QAction("查找上一个", self)
        find_prev_action.setShortcut(QKeySequence.FindPrevious)
        find_prev_action.triggered.connect(self.search_prev)
        edit_menu.addAction(find_prev_action)

        # 视图菜单
        view_menu = menubar.addMenu("👁️ 视图")

        zoom_in_action = QAction("放大", self)
        zoom_in_action.setShortcut(QKeySequence.ZoomIn)
        zoom_in_action.triggered.connect(self.zoom_in)
        view_menu.addAction(zoom_in_action)

        zoom_out_action = QAction("缩小", self)
        zoom_out_action.setShortcut(QKeySequence.ZoomOut)
        zoom_out_action.triggered.connect(self.zoom_out)
        view_menu.addAction(zoom_out_action)

        zoom_reset_action = QAction("实际大小", self)
        zoom_reset_action.setShortcut("Ctrl+0")
        zoom_reset_action.triggered.connect(self.zoom_reset)
        view_menu.addAction(zoom_reset_action)

        view_menu.addSeparator()

        fit_width_action = QAction("适合宽度", self)
        fit_width_action.triggered.connect(self.fit_width)
        view_menu.addAction(fit_width_action)

        fit_page_action = QAction("适合页面", self)
        fit_page_action.triggered.connect(self.fit_page)
        view_menu.addAction(fit_page_action)

        view_menu.addSeparator()

        rotate_action = QAction("旋转页面", self)
        rotate_action.setShortcut("Ctrl+R")
        rotate_action.triggered.connect(self.rotate_page)
        view_menu.addAction(rotate_action)

        view_menu.addSeparator()

        theme_menu = view_menu.addMenu("🎨 主题")

        light_action = QAction("亮色主题", self)
        light_action.triggered.connect(lambda: self.change_theme("light"))
        theme_menu.addAction(light_action)

        dark_action = QAction("暗色主题", self)
        dark_action.triggered.connect(lambda: self.change_theme("dark"))
        theme_menu.addAction(dark_action)

        eye_care_action = QAction("护眼模式", self)
        eye_care_action.triggered.connect(lambda: self.change_theme("eye_care"))
        theme_menu.addAction(eye_care_action)

        view_menu.addSeparator()

        fullscreen_action = QAction("全屏模式", self)
        fullscreen_action.setShortcut(Qt.Key_F11)
        fullscreen_action.triggered.connect(self.toggle_fullscreen)
        view_menu.addAction(fullscreen_action)

        # 导航菜单
        nav_menu = menubar.addMenu("🧭 导航")

        prev_page_action = QAction("上一页", self)
        prev_page_action.setShortcut(Qt.Key_PageUp)
        prev_page_action.triggered.connect(self.previous_page)
        nav_menu.addAction(prev_page_action)

        next_page_action = QAction("下一页", self)
        next_page_action.setShortcut(Qt.Key_PageDown)
        next_page_action.triggered.connect(self.next_page)
        nav_menu.addAction(next_page_action)

        nav_menu.addSeparator()

        go_to_page_action = QAction("跳转到页面...", self)
        go_to_page_action.setShortcut("Ctrl+G")
        go_to_page_action.triggered.connect(self._show_goto_dialog)
        nav_menu.addAction(go_to_page_action)

        nav_menu.addSeparator()

        back_action = QAction("后退", self)
        back_action.setShortcut(Qt.ALT + Qt.Key_Left)
        back_action.triggered.connect(self.go_back)
        nav_menu.addAction(back_action)

        forward_action = QAction("前进", self)
        forward_action.setShortcut(Qt.ALT + Qt.Key_Right)
        forward_action.triggered.connect(self.go_forward)
        nav_menu.addAction(forward_action)

        # 工具菜单
        tools_menu = menubar.addMenu("🔧 工具")

        bookmark_action = QAction("添加书签", self)
        bookmark_action.setShortcut("Ctrl+D")
        bookmark_action.triggered.connect(self.add_bookmark)
        tools_menu.addAction(bookmark_action)

        note_action = QAction("添加笔记", self)
        note_action.setShortcut("Ctrl+N")
        note_action.triggered.connect(self.add_note)
        tools_menu.addAction(note_action)

        tools_menu.addSeparator()

        if OCR_AVAILABLE:
            ocr_action = QAction("OCR识别当前页", self)
            ocr_action.triggered.connect(self.ocr_current_page)
            tools_menu.addAction(ocr_action)

            # 帮助菜单
        help_menu = menubar.addMenu("❓ 帮助")

        shortcuts_action = QAction("快捷键说明", self)
        shortcuts_action.setShortcut(Qt.Key_F1)
        shortcuts_action.triggered.connect(self._show_shortcuts_help)
        help_menu.addAction(shortcuts_action)

        about_action = QAction("关于", self)
        about_action.triggered.connect(self._show_about)
        help_menu.addAction(about_action)

    def _build_toolbar(self):
        """构建工具栏"""
        toolbar = QToolBar("主工具栏")
        toolbar.setMovable(False)
        toolbar.setIconSize(QSize(24, 24))
        toolbar.setToolButtonStyle(Qt.ToolButtonTextBesideIcon)
        self.addToolBar(toolbar)

        # 文件操作
        open_btn = QPushButton("📂 打开")
        open_btn.setToolTip("打开文档 (Ctrl+O)")
        open_btn.clicked.connect(self.open_file)
        toolbar.addWidget(open_btn)

        save_btn = QPushButton("💾 保存")
        save_btn.setToolTip("保存文档 (Ctrl+S)")
        save_btn.clicked.connect(self.save_current_document)
        toolbar.addWidget(save_btn)

        toolbar.addSeparator()

        # 缩放控制
        zoom_out_btn = QPushButton("🔍-")
        zoom_out_btn.setToolTip("缩小 (Ctrl+-)")
        zoom_out_btn.clicked.connect(self.zoom_out)
        toolbar.addWidget(zoom_out_btn)

        self.zoom_label = QLabel("100%")
        self.zoom_label.setMinimumWidth(60)
        self.zoom_label.setAlignment(Qt.AlignCenter)
        self.zoom_label.setStyleSheet("font-weight: bold;")
        toolbar.addWidget(self.zoom_label)

        zoom_in_btn = QPushButton("🔍+")
        zoom_in_btn.setToolTip("放大 (Ctrl++)")
        zoom_in_btn.clicked.connect(self.zoom_in)
        toolbar.addWidget(zoom_in_btn)

        self.zoom_combo = QComboBox()
        self.zoom_combo.addItems([
            "25%", "50%", "75%", "100%", "125%",
            "150%", "200%", "300%", "适合宽度", "适合页面"
        ])
        self.zoom_combo.setCurrentText("100%")
        self.zoom_combo.currentTextChanged.connect(self._on_zoom_preset)
        toolbar.addWidget(self.zoom_combo)

        toolbar.addSeparator()

        # 页面导航
        back_btn = QPushButton("◀◀")
        back_btn.setToolTip("后退 (Alt+Left)")
        back_btn.clicked.connect(self.go_back)
        toolbar.addWidget(back_btn)

        prev_btn = QPushButton("◀")
        prev_btn.setToolTip("上一页 (PageUp)")
        prev_btn.clicked.connect(self.previous_page)
        toolbar.addWidget(prev_btn)

        self.page_input = QLineEdit()
        self.page_input.setMaximumWidth(60)
        self.page_input.setAlignment(Qt.AlignCenter)
        self.page_input.setPlaceholderText("页码")
        self.page_input.returnPressed.connect(self._jump_to_input_page)
        toolbar.addWidget(self.page_input)

        self.total_pages_label = QLabel("/ 0")
        toolbar.addWidget(self.total_pages_label)

        next_btn = QPushButton("▶")
        next_btn.setToolTip("下一页 (PageDown)")
        next_btn.clicked.connect(self.next_page)
        toolbar.addWidget(next_btn)

        forward_btn = QPushButton("▶▶")
        forward_btn.setToolTip("前进 (Alt+Right)")
        forward_btn.clicked.connect(self.go_forward)
        toolbar.addWidget(forward_btn)

        toolbar.addSeparator()

        # 视图控制
        rotate_btn = QPushButton("🔄")
        rotate_btn.setToolTip("旋转页面 (Ctrl+R)")
        rotate_btn.clicked.connect(self.rotate_page)
        toolbar.addWidget(rotate_btn)

        self.theme_btn = QPushButton("🌙")
        self.theme_btn.setToolTip("切换主题")
        self.theme_btn.clicked.connect(self._cycle_theme)
        toolbar.addWidget(self.theme_btn)

        toolbar.addSeparator()

        # 书签和笔记
        bookmark_btn = QPushButton("⭐")
        bookmark_btn.setToolTip("添加书签 (Ctrl+D)")
        bookmark_btn.clicked.connect(self.add_bookmark)
        toolbar.addWidget(bookmark_btn)

        note_btn = QPushButton("📝")
        note_btn.setToolTip("添加笔记 (Ctrl+N)")
        note_btn.clicked.connect(self.add_note)
        toolbar.addWidget(note_btn)

        toolbar.addSeparator()

        # 搜索
        self.search_input = QLineEdit()
        self.search_input.setPlaceholderText("🔍 搜索...")
        self.search_input.setMaximumWidth(200)
        self.search_input.returnPressed.connect(self.search_text)
        toolbar.addWidget(self.search_input)

        search_prev_btn = QPushButton("⏮")
        search_prev_btn.setToolTip("上一个匹配")
        search_prev_btn.clicked.connect(self.search_prev)
        toolbar.addWidget(search_prev_btn)

        search_next_btn = QPushButton("⏭")
        search_next_btn.setToolTip("下一个匹配")
        search_next_btn.clicked.connect(self.search_next)
        toolbar.addWidget(search_next_btn)

        toolbar.addSeparator()

        # 注释工具
        self.highlight_btn = QPushButton("🖍")
        self.highlight_btn.setCheckable(True)
        self.highlight_btn.setToolTip("高亮模式")
        self.highlight_btn.clicked.connect(lambda: self._toggle_annotation_mode("highlight"))
        toolbar.addWidget(self.highlight_btn)

        self.underline_btn = QPushButton("〰")
        self.underline_btn.setCheckable(True)
        self.underline_btn.setToolTip("下划线模式")
        self.underline_btn.clicked.connect(lambda: self._toggle_annotation_mode("underline"))
        toolbar.addWidget(self.underline_btn)

        self.note_ann_btn = QPushButton("💬")
        self.note_ann_btn.setCheckable(True)
        self.note_ann_btn.setToolTip("批注模式")
        self.note_ann_btn.clicked.connect(lambda: self._toggle_annotation_mode("note"))
        toolbar.addWidget(self.note_ann_btn)

    def _build_statusbar(self):
        """构建状态栏"""
        self.status_bar = QStatusBar()
        self.setStatusBar(self.status_bar)

        # 文档格式
        self.format_label = QLabel("就绪")
        self.format_label.setStyleSheet("padding: 0 10px;")
        self.status_bar.addWidget(self.format_label)

        self.status_bar.addWidget(QLabel("|"))

        # 页面信息
        self.page_info_label = QLabel("无文档")
        self.page_info_label.setMinimumWidth(150)
        self.status_bar.addWidget(self.page_info_label)

        self.status_bar.addWidget(QLabel("|"))

        # 阅读进度
        self.progress_label = QLabel("进度: 0%")
        self.progress_label.setMinimumWidth(100)
        self.status_bar.addWidget(self.progress_label)

        self.status_bar.addWidget(QLabel("|"))

        # 阅读时间
        self.time_label = QLabel("阅读时间: 00:00:00")
        self.time_label.setMinimumWidth(180)
        self.status_bar.addPermanentWidget(self.time_label)

    def _setup_shortcuts(self):
        """设置快捷键"""
        # 已在菜单栏中设置
        pass

        # ==================== 文件操作 ====================

    def open_file(self):
        """打开文件"""
        config = ConfigManager.load()
        last_dir = config.get("last_opened_dir", str(Path.cwd()))

        # 构建文件过滤器
        filters = []

        # PyMuPDF
        filters.append("PDF文档 (*.pdf)")
        filters.append("EPUB电子书 (*.epub)")
        filters.append("MOBI电子书 (*.mobi *.azw *.azw3)")
        filters.append("FictionBook (*.fb2)")
        filters.append("XPS文档 (*.xps)")
        filters.append("漫画书 (*.cbz *.cbr)")

        # 文本
        filters.append("文本文件 (*.txt)")
        filters.append("Markdown (*.md *.markdown)")

        if DOCX_AVAILABLE:
            filters.append("Word文档 (*.docx)")

        if ODT_AVAILABLE:
            filters.append("OpenDocument (*.odt *.odp *.ods)")

        if PPTX_AVAILABLE:
            filters.append("PowerPoint (*.pptx)")

            # 所有支持的格式
        all_exts = "*.pdf *.epub *.mobi *.azw *.azw3 *.fb2 *.xps *.cbz *.cbr *.txt *.md *.markdown"
        if DOCX_AVAILABLE:
            all_exts += " *.docx"
        if ODT_AVAILABLE:
            all_exts += " *.odt *.odp *.ods"
        if PPTX_AVAILABLE:
            all_exts += " *.pptx"

        filters.insert(0, f"所有支持的格式 ({all_exts})")

        filter_str = ";;".join(filters)

        file_path, _ = QFileDialog.getOpenFileName(
            self, "打开文档", last_dir, filter_str
        )

        if not file_path:
            return

            # 保存最后打开的目录
        config["last_opened_dir"] = str(Path(file_path).parent)
        ConfigManager.save(config)

        # 创建新视图
        self._open_document_in_new_tab(file_path)

    def _open_document_in_new_tab(self, file_path: str):
        """在新标签页中打开文档"""
        try:
            view = DocumentView(self)

            # 连接信号
            view.document_loaded.connect(self._on_document_loaded)
            view.page_changed.connect(self._on_page_changed)
            view.zoom_changed.connect(self._on_zoom_changed)

            # 加载文档
            view.load_document(Path(file_path))

        except Exception as e:
            QMessageBox.critical(self, "打开失败", f"无法打开文件:\n{str(e)}")

    @pyqtSlot(str, str, int)
    def _on_document_loaded(self, file_path: str, fmt: str, total_pages: int):
        """文档加载完成"""
        view = self.sender()

        if not isinstance(view, DocumentView):
            return

            # 添加到标签页
        title = Path(file_path).name
        if len(title) > 30:
            title = title[:27] + "..."

        idx = self.tab_widget.indexOf(view)
        if idx < 0:
            self.tab_widget.addTab(view, title)
            self.tab_widget.setCurrentWidget(view)
        else:
            self.tab_widget.setTabText(idx, title)

            # 更新UI
        self.format_label.setText(f"{fmt.upper()} | {total_pages}页")
        self._update_page_display(1, total_pages)

        # 添加到最近文件
        self._add_to_recent_files(file_path)

        self.show_status(f"已打开: {Path(file_path).name}", 4000)

    def save_current_document(self):
        """保存当前文档"""
        view = self._current_view()
        if not view:
            return

        if not view.is_text_editable:
            QMessageBox.information(
                self, "保存",
                "当前文档类型不支持直接保存编辑。\n您可以使用\"导出\"功能。"
            )
            return

        if not view.current_file:
            save_path, _ = QFileDialog.getSaveFileName(
                self, "另存为", "",
                "文本文件 (*.txt);;Markdown (*.md);;所有文件 (*)"
            )
            if not save_path:
                return
            view.current_file = save_path

            # 收集所有文本内容
        content = []
        for widget in view.page_widgets:
            if isinstance(widget, TextPageWidget):
                content.append(widget.toPlainText())

        text = "\n\n".join(content)

        try:
            with open(view.current_file, "w", encoding="utf-8") as f:
                f.write(text)

            self.show_status("保存成功", 3000)
        except Exception as e:
            QMessageBox.critical(self, "保存失败", str(e))

    def export_pages(self):
        """导出页面"""
        view = self._current_view()
        if view:
            view.export_pages()

    def print_preview(self):
        """打印预览"""
        view = self._current_view()
        if view:
            view.print_preview()

    def print_document(self):
        """打印文档"""
        view = self._current_view()
        if view:
            view.print_document()

            # ==================== 标签页管理 ====================

    def _close_tab(self, index: int):
        """关闭标签页"""
        widget = self.tab_widget.widget(index)

        if isinstance(widget, DocumentView):
            widget.close_view()

        self.tab_widget.removeTab(index)

        if self.tab_widget.count() == 0:
            self._reset_ui()

    def _detach_tab(self, index: int):
        """分离标签页为独立窗口"""
        if index < 0:
            return

        widget = self.tab_widget.widget(index)
        if not isinstance(widget, DocumentView):
            return

            # 从当前窗口移除
        self.tab_widget.removeTab(index)

        # 创建新窗口
        new_window = UniversalDocumentReader()
        new_window.show()

        # 重新绑定视图
        widget.setParent(new_window)
        widget.main = new_window

        # 添加到新窗口
        title = widget.current_file and Path(widget.current_file).name or "未命名"
        new_window.tab_widget.addTab(widget, title)
        new_window.tab_widget.setCurrentWidget(widget)

        # 保持引用防止GC
        self.detached_windows.append(new_window)

    def _on_tab_changed(self, index: int):
        """标签页切换"""
        self._update_reading_info()

        view = self._current_view()
        if view:
            info = view.get_progress_info()
            self._update_page_display(
                info["current_page"],
                info["total_pages"]
            )
            self.format_label.setText(
                f"{view.current_format.upper() if view.current_format else ''} | "
                f"{info['total_pages']}页"
            )
        else:
            self._reset_ui()

    def _reset_ui(self):
        """重置UI（无文档时）"""
        self.format_label.setText("就绪")
        self.page_info_label.setText("无文档")
        self.progress_label.setText("进度: 0%")
        self.time_label.setText("阅读时间: 00:00:00")
        self.total_pages_label.setText("/ 0")
        self.page_input.clear()
        self.zoom_label.setText("100%")

        # ==================== 视图控制 ====================

    def zoom_in(self):
        """放大"""
        view = self._current_view()
        if view:
            view.zoom_in()

    def zoom_out(self):
        """缩小"""
        view = self._current_view()
        if view:
            view.zoom_out()

    def zoom_reset(self):
        """重置缩放"""
        view = self._current_view()
        if view:
            view.set_zoom(100)

    def fit_width(self):
        """适合宽度"""
        view = self._current_view()
        if view:
            view.fit_width()

    def fit_page(self):
        """适合页面"""
        view = self._current_view()
        if view:
            view.fit_page()

    def _on_zoom_preset(self, text: str):
        """缩放预设值改变"""
        view = self._current_view()
        if not view:
            return

        if text == "适合宽度":
            view.fit_width()
        elif text == "适合页面":
            view.fit_page()
        else:
            try:
                zoom = int(text.replace("%", ""))
                view.set_zoom(zoom)
            except:
                pass

    @pyqtSlot(int)
    def _on_zoom_changed(self, zoom: int):
        """缩放级别改变"""
        self.zoom_label.setText(f"{zoom}%")

        # 更新下拉框
        zoom_text = f"{zoom}%"
        if zoom_text in [self.zoom_combo.itemText(i) for i in range(self.zoom_combo.count())]:
            self.zoom_combo.blockSignals(True)
            self.zoom_combo.setCurrentText(zoom_text)
            self.zoom_combo.blockSignals(False)

    def rotate_page(self):
        """旋转页面"""
        view = self._current_view()
        if view:
            view.rotate_page()

    def change_theme(self, theme: str):
        """更改主题"""
        self.current_theme = theme

        # 更新所有视图
        for i in range(self.tab_widget.count()):
            widget = self.tab_widget.widget(i)
            if isinstance(widget, DocumentView):
                widget.toggle_theme(theme)

                # 更新主窗口样式
        self._apply_theme()

        # 更新按钮图标
        theme_icons = {
            "light": "🌙",
            "dark": "☀️",
            "eye_care": "🌿"
        }
        self.theme_btn.setText(theme_icons.get(theme, "🌙"))

        # 保存配置
        config = ConfigManager.load()
        config["theme"] = theme
        ConfigManager.save(config)

    def _cycle_theme(self):
        """循环切换主题"""
        themes = ["light", "dark", "eye_care"]
        current_index = themes.index(self.current_theme)
        next_theme = themes[(current_index + 1) % len(themes)]
        self.change_theme(next_theme)

    def _apply_theme(self):
        """应用主题样式"""
        if self.current_theme == "dark":
            self.setStyleSheet("""  
                QMainWindow {  
                    background-color: #1E1E1E;  
                    color: #D4D4D4;  
                }  
                QMenuBar {  
                    background-color: #2D2D30;  
                    color: #D4D4D4;  
                }  
                QMenuBar::item:selected {  
                    background-color: #3E3E40;  
                }  
                QMenu {  
                    background-color: #2D2D30;  
                    color: #D4D4D4;  
                }  
                QMenu::item:selected {  
                    background-color: #3E3E40;  
                }  
                QToolBar {  
                    background-color: #2D2D30;  
                    border-bottom: 1px solid #3E3E40;  
                }  
                QPushButton {  
                    background-color: #3E3E40;  
                    color: #D4D4D4;  
                    border: none;  
                    padding: 5px 10px;  
                    border-radius: 3px;  
                }  
                QPushButton:hover {  
                    background-color: #4E4E50;  
                }  
                QLineEdit, QComboBox {  
                    background-color: #3E3E40;  
                    color: #D4D4D4;  
                    border: 1px solid #555555;  
                    padding: 3px;  
                }  
                QTabWidget::pane {  
                    border: 1px solid #3E3E40;  
                }  
                QTabBar::tab {  
                    background-color: #2D2D30;  
                    color: #D4D4D4;  
                    padding: 8px 16px;  
                    border: 1px solid #3E3E40;  
                }  
                QTabBar::tab:selected {  
                    background-color: #3E3E40;  
                }  
                QStatusBar {  
                    background-color: #2D2D30;  
                    color: #D4D4D4;  
                }  
            """)
        elif self.current_theme == "eye_care":
            self.setStyleSheet("""  
                QMainWindow {  
                    background-color: #C7EDCC;  
                    color: #2F4F2F;  
                }  
                QMenuBar {  
                    background-color: #A8D5BA;  
                    color: #2F4F2F;  
                }  
                QToolBar {  
                    background-color: #A8D5BA;  
                    border-bottom: 1px solid #8FBC8F;  
                }  
                QPushButton {  
                    background-color: #8FBC8F;  
                    color: #2F4F2F;  
                    border: none;  
                    padding: 5px 10px;  
                    border-radius: 3px;  
                }  
                QPushButton:hover {  
                    background-color: #7CAC7C;  
                }  
                QLineEdit, QComboBox {  
                    background-color: white;  
                    color: #2F4F2F;  
                    border: 1px solid #8FBC8F;  
                    padding: 3px;  
                }  
            """)
        else:
            self.setStyleSheet("")

    def toggle_fullscreen(self):
        """切换全屏"""
        if self.isFullScreen():
            self.showNormal()
        else:
            self.showFullScreen()

            # ==================== 导航 ====================

    def previous_page(self):
        """上一页"""
        view = self._current_view()
        if view:
            view.previous_page()

    def next_page(self):
        """下一页"""
        view = self._current_view()
        if view:
            view.next_page()

    def _jump_to_input_page(self):
        """跳转到输入的页码"""
        view = self._current_view()
        if not view:
            return

        try:
            page = int(self.page_input.text())
            if 1 <= page <= len(view.page_widgets):
                view.jump_to_page(page - 1)
            else:
                self.show_status("页码超出范围", 2000)
        except ValueError:
            self.show_status("请输入有效的页码", 2000)

    def _show_goto_dialog(self):
        """显示跳转对话框"""
        view = self._current_view()
        if not view:
            return

        page, ok = QInputDialog.getInt(
            self, "跳转到页面",
            f"输入页码 (1-{len(view.page_widgets)}):",
            value=view.current_page_index + 1,
            min=1, max=len(view.page_widgets)
        )

        if ok:
            view.jump_to_page(page - 1)

    def go_back(self):
        """后退"""
        view = self._current_view()
        if view:
            view.go_back()

    def go_forward(self):
        """前进"""
        view = self._current_view()
        if view:
            view.go_forward()

    @pyqtSlot(int, int)
    def _on_page_changed(self, current: int, total: int):
        """页面改变"""
        self._update_page_display(current, total)

    def _update_page_display(self, current: int, total: int):
        """更新页面显示"""
        self.page_input.setText(str(current))
        self.total_pages_label.setText(f"/ {total}")
        self.page_info_label.setText(f"第 {current} / {total} 页")

        if total > 0:
            progress = (current / total) * 100
            self.progress_label.setText(f"进度: {progress:.1f}%")

            # ==================== 书签和笔记 ====================

    def add_bookmark(self):
        """添加书签"""
        view = self._current_view()
        if view:
            view.add_bookmark()

    def add_note(self):
        """添加笔记"""
        view = self._current_view()
        if view:
            view.add_note()

            # ==================== 注释 ====================

    def _toggle_annotation_mode(self, mode: str):
        """切换注释模式"""
        view = self._current_view()
        if not view:
            return

            # 互斥按钮
        btns = {
            "highlight": self.highlight_btn,
            "underline": self.underline_btn,
            "note": self.note_ann_btn
        }

        current_btn = btns[mode]

        # 如果当前按钮已选中，则取消
        if current_btn.isChecked():
            view.set_annotation_mode(mode)
            # 取消其他按钮
            for m, btn in btns.items():
                if m != mode:
                    btn.setChecked(False)
        else:
            view.set_annotation_mode(None)

            # ==================== 搜索 ====================

    def search_text(self):
        """搜索文本"""
        view = self._current_view()
        if not view:
            return

        term = self.search_input.text().strip()
        if not term:
            self.show_status("请输入搜索内容", 2000)
            return

        view.search_text(term)

    def search_next(self):
        """下一个搜索结果"""
        view = self._current_view()
        if view:
            view.search_next()

    def search_prev(self):
        """上一个搜索结果"""
        view = self._current_view()
        if view:
            view.search_prev()

            # ==================== OCR ====================

    def ocr_current_page(self):
        """OCR识别当前页"""
        view = self._current_view()
        if view:
            view.ocr_current_page()

            # ==================== 最近文件 ====================

    def _add_to_recent_files(self, file_path: str):
        """添加到最近文件"""
        config = ConfigManager.load()
        recent = config.get("recent_files", [])

        if file_path in recent:
            recent.remove(file_path)

        recent.insert(0, file_path)
        config["recent_files"] = recent[:50]  # 保留最近50个

        ConfigManager.save(config)

    def _populate_recent_menu(self, menu: QMenu):
        """填充最近文件菜单"""
        menu.clear()

        config = ConfigManager.load()
        recent = config.get("recent_files", [])

        if not recent:
            action = QAction("(无最近文件)", self)
            action.setEnabled(False)
            menu.addAction(action)
            return

        for file_path in recent[:20]:  # 显示最近20个
            if Path(file_path).exists():
                name = Path(file_path).name
                action = QAction(name, self)
                action.setToolTip(file_path)
                action.triggered.connect(
                    lambda checked, p=file_path: self._open_document_in_new_tab(p)
                )
                menu.addAction(action)

        menu.addSeparator()

        clear_action = QAction("清除历史记录", self)
        clear_action.triggered.connect(self._clear_recent_files)
        menu.addAction(clear_action)

    def _clear_recent_files(self):
        """清除最近文件"""
        reply = QMessageBox.question(
            self, "确认",
            "确定要清除所有最近文件记录吗？",
            QMessageBox.Yes | QMessageBox.No
        )

        if reply == QMessageBox.Yes:
            config = ConfigManager.load()
            config["recent_files"] = []
            ConfigManager.save(config)
            self.show_status("已清除历史记录", 2000)

            # ==================== 帮助 ====================

    def _show_shortcuts_help(self):
        """显示快捷键帮助"""
        help_text = """  
        <h2>快捷键说明</h2>  

        <h3>📂 文件操作</h3>  
        <ul>  
            <li><b>Ctrl+O</b> - 打开文档</li>  
            <li><b>Ctrl+S</b> - 保存文档</li>  
            <li><b>Ctrl+P</b> - 打印</li>  
            <li><b>Ctrl+W</b> - 关闭标签页</li>  
            <li><b>Ctrl+Q</b> - 退出</li>  
        </ul>  

        <h3>🔍 查找</h3>  
        <ul>  
            <li><b>Ctrl+F</b> - 查找</li>  
                        <li><b>F3</b> - 查找下一个</li>
            <li><b>Shift+F3</b> - 查找上一个</li>
        </ul>
        
        <h3>👁️ 视图</h3>
        <ul>
            <li><b>Ctrl+加号</b> - 放大</li>
            <li><b>Ctrl+减号</b> - 缩小</li>
            <li><b>Ctrl+0</b> - 实际大小</li>
            <li><b>Ctrl+R</b> - 旋转页面</li>
            <li><b>F11</b> - 全屏模式</li>
        </ul>
        
        <h3>🧭 导航</h3>
        <ul>
            <li><b>PageUp</b> - 上一页</li>
            <li><b>PageDown</b> - 下一页</li>
            <li><b>Home</b> - 第一页</li>
            <li><b>End</b> - 最后一页</li>
            <li><b>Ctrl+G</b> - 跳转到页面</li>
            <li><b>Alt+Left</b> - 后退</li>
            <li><b>Alt+Right</b> - 前进</li>
        </ul>
        
        <h3>🔧 工具</h3>
        <ul>
            <li><b>Ctrl+D</b> - 添加书签</li>
            <li><b>Ctrl+N</b> - 添加笔记</li>
        </ul>
        
        <h3>❓ 帮助</h3>
        <ul>
            <li><b>F1</b> - 快捷键说明</li>
        </ul>
        """

        dialog = QDialog(self)
        dialog.setWindowTitle("快捷键说明")
        dialog.setMinimumSize(500, 600)

        layout = QVBoxLayout(dialog)

        browser = QTextBrowser()
        browser.setHtml(help_text)
        browser.setOpenExternalLinks(True)
        layout.addWidget(browser)

        close_btn = QPushButton("关闭")
        close_btn.clicked.connect(dialog.close)
        layout.addWidget(close_btn)

        dialog.exec_()

    def _show_about(self):
        """显示关于对话框"""
        about_text = f"""
        <h2>专业文档阅读器</h2>
        <p><b>版本:</b> 2.0.0 (PyQt5重构版)</p>
        
        <h3>✨ 主要特性</h3>
        <ul>
            <li>✅ 支持多种格式: PDF, EPUB, MOBI, DOCX, TXT, Markdown等</li>
            <li>✅ 多标签页浏览</li>
            <li>✅ 智能书签与注释系统</li>
            <li>✅ 全文搜索与高亮</li>
            <li>✅ 阅读进度跟踪</li>
            <li>✅ OCR文字识别</li>
            <li>✅ 多主题支持（亮色/暗色/护眼）</li>
            <li>✅ 打印与导出功能</li>
            <li>✅ 自动保存阅读位置</li>
        </ul>
        
        <h3>📦 依赖库</h3>
        <ul>
            <li>PyQt5 - GUI框架</li>
            <li>PyMuPDF (fitz) - PDF渲染核心</li>
            <li>python-docx - Word文档支持 {'✅' if DOCX_AVAILABLE else '❌'}</li>
            <li>odfpy - OpenDocument支持 {'✅' if ODT_AVAILABLE else '❌'}</li>
            <li>python-pptx - PowerPoint支持 {'✅' if PPTX_AVAILABLE else '❌'}</li>
            <li>pytesseract - OCR识别 {'✅' if OCR_AVAILABLE else '❌'}</li>
        </ul>
        
        <p style="margin-top: 20px;">
            <b>作者:</b> UltimateReader Team<br>
            <b>许可:</b> MIT License
        </p>
        """

        QMessageBox.about(self, "关于", about_text)

    # ==================== 状态更新 ====================

    def _update_reading_info(self):
        """更新阅读信息"""
        view = self._current_view()
        if not view:
            return

        info = view.get_progress_info()

        # 更新进度
        self.progress_label.setText(f"进度: {info['progress_percent']:.1f}%")

        # 更新时间
        self.time_label.setText(f"阅读时间: {info['time_formatted']}")

    def show_status(self, message: str, timeout: int = 5000):
        """显示状态栏消息"""
        self.status_bar.showMessage(message, timeout)

    # ==================== 辅助方法 ====================

    def _current_view(self) -> Optional[DocumentView]:
        """获取当前视图"""
        widget = self.tab_widget.currentWidget()
        if isinstance(widget, DocumentView):
            return widget
        return None

    # ==================== 拖放支持 ====================

    def dragEnterEvent(self, event: QDragEnterEvent):
        """拖入事件"""
        if event.mimeData().hasUrls():
            event.acceptProposedAction()

    def dropEvent(self, event: QDropEvent):
        """放下事件"""
        for url in event.mimeData().urls():
            file_path = url.toLocalFile()
            if Path(file_path).is_file():
                try:
                    self._open_document_in_new_tab(file_path)
                except Exception as e:
                    QMessageBox.warning(
                        self, "打开失败",
                        f"无法打开文件 {Path(file_path).name}:\n{str(e)}"
                    )

    # ==================== 窗口事件 ====================

    def _load_config(self):
        """加载配置"""
        config = ConfigManager.load()

        # 窗口位置和大小
        window = config.get("window", {})
        self.setGeometry(
            window.get("x", 100),
            window.get("y", 100),
            window.get("width", 1600),
            window.get("height", 900)
        )

        # 主题
        theme = config.get("theme", "light")
        self.change_theme(theme)

    def closeEvent(self, event):
        """关闭事件"""
        # 保存所有文档状态
        for i in range(self.tab_widget.count()):
            widget = self.tab_widget.widget(i)
            if isinstance(widget, DocumentView):
                widget.close_view()

        # 保存窗口配置
        config = ConfigManager.load()
        config["window"] = {
            "x": self.x(),
            "y": self.y(),
            "width": self.width(),
            "height": self.height()
        }
        ConfigManager.save(config)

        event.accept()


# ==================== 主程序入口 ====================
def main():
    """主函数"""
    import sys

    # 启用高DPI支持
    QApplication.setAttribute(Qt.AA_EnableHighDpiScaling, True)
    QApplication.setAttribute(Qt.AA_UseHighDpiPixmaps, True)

    app = QApplication(sys.argv)
    app.setApplicationName("专业文档阅读器")
    app.setOrganizationName("UltimateReader")

    # 设置应用图标（如果有）
    # app.setWindowIcon(QIcon("icon.png"))

    window = UniversalDocumentReader()
    window.show()

    # 如果有命令行参数，打开文件
    if len(sys.argv) > 1:
        for file_path in sys.argv[1:]:
            if Path(file_path).exists():
                try:
                    window._open_document_in_new_tab(file_path)
                except Exception as e:
                    print(f"无法打开文件 {file_path}: {e}")

    sys.exit(app.exec_())


if __name__ == "__main__":
    main()